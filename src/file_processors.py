import os
import io
import PyPDF2
from docx import Document
import pandas as pd
from pptx import Presentation


def extract_text_from_pdf(file_like):
    """
    Extracts text content from PDF bytes.
    
    Args:
        file_like: File-like object containing PDF bytes
    
    Returns:
        str: Extracted text content
    """
    try:
        pdf_reader = PyPDF2.PdfReader(file_like)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from PDF bytes: {e}")
        return ""


def extract_text_from_docx(file_like):
    """
    Extracts text content from DOCX bytes.
    
    Args:
        file_like: File-like object containing DOCX bytes
    
    Returns:
        str: Extracted text content
    """
    try:
        doc = Document(file_like)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from DOCX bytes: {e}")
        return ""


def extract_text_from_csv(file_like):
    """
    Extracts text content from CSV bytes.
    
    Args:
        file_like: File-like object containing CSV bytes
    
    Returns:
        str: Extracted text content
    """
    try:
        df = pd.read_csv(file_like)
        text = df.to_string(index=False)
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from CSV bytes: {e}")
        return ""


def extract_text_from_pptx(file_like):
    """
    Extracts text content from PPTX bytes.
    
    Args:
        file_like: File-like object containing PPTX bytes
    
    Returns:
        str: Extracted text content
    """
    try:
        prs = Presentation(file_like)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from PPTX bytes: {e}")
        return ""


def extract_text_from_xlsx(file_like):
    """
    Extracts text content from XLSX bytes.
    
    Args:
        file_like: File-like object containing XLSX bytes
    
    Returns:
        str: Extracted text content
    """
    try:
        df = pd.read_excel(file_like)
        text = df.to_string(index=False)
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from XLSX bytes: {e}")
        return ""



def extract_text_from_pdf_file(filepath):
    """
    Extracts text content from a PDF file.
    
    Args:
        filepath (str): Path to the PDF file
    
    Returns:
        str: Extracted text content
    """
    try:
        with open(filepath, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from PDF {filepath}: {e}")
        return ""


def extract_text_from_docx_file(filepath):
    """
    Extracts text content from a DOCX file.
    
    Args:
        filepath (str): Path to the DOCX file
    
    Returns:
        str: Extracted text content
    """
    try:
        doc = Document(filepath)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from DOCX {filepath}: {e}")
        return ""


def extract_text_from_csv_file(filepath):
    """
    Extracts text content from a CSV file.
    
    Args:
        filepath (str): Path to the CSV file
    
    Returns:
        str: Extracted text content
    """
    try:
        df = pd.read_csv(filepath)
        text = df.to_string(index=False)
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from CSV {filepath}: {e}")
        return ""


def extract_text_from_pptx_file(filepath):
    """
    Extracts text content from a PPTX file.
    
    Args:
        filepath (str): Path to the PPTX file
    
    Returns:
        str: Extracted text content
    """
    try:
        prs = Presentation(filepath)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from PPTX {filepath}: {e}")
        return ""


def extract_text_from_xlsx_file(filepath):
    """
    Extracts text content from an XLSX file.
    
    Args:
        filepath (str): Path to the XLSX file
    
    Returns:
        str: Extracted text content
    """
    try:
        df = pd.read_excel(filepath)
        text = df.to_string(index=False)
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from XLSX {filepath}: {e}")
        return ""


def extract_text_from_file(filepath):
    """
    Extracts text content from various file formats.
    
    Args:
        filepath (str): Path to the file
    
    Returns:
        str: Extracted text content
    """
    filename = os.path.basename(filepath)
    file_extension = os.path.splitext(filename)[1].lower()
    
    if file_extension == '.txt':
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            print(f"Error reading TXT file {filepath}: {e}")
            return ""
    elif file_extension == '.pdf':
        return extract_text_from_pdf_file(filepath)
    elif file_extension == '.docx':
        return extract_text_from_docx_file(filepath)
    elif file_extension == '.csv':
        return extract_text_from_csv_file(filepath)
    elif file_extension == '.pptx':
        return extract_text_from_pptx_file(filepath)
    elif file_extension in ['.xlsx', '.xls']:
        return extract_text_from_xlsx_file(filepath)
    else:
        print(f"Unsupported file format: {file_extension}")
        return ""